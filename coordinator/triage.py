import logging
import re
from urllib.parse import urlparse

log = logging.getLogger(__name__)

YOUTUBE_DOMAINS = {"youtube.com", "youtu.be", "www.youtube.com", "m.youtube.com"}
TWITTER_DOMAINS = {"twitter.com", "x.com", "www.twitter.com"}
REDDIT_DOMAINS = {"reddit.com", "www.reddit.com", "old.reddit.com"}
ARXIV_DOMAINS = {"arxiv.org", "www.arxiv.org"}

SUBTITLES_MIN_CHARS = 500
SUBTITLES_MIN_UNIQUE_RATIO = 0.6


def detect_url_type(url: str) -> str:
    try:
        domain = urlparse(url).netloc.lower().lstrip("www.")
    except Exception:
        return "article"

    if domain in YOUTUBE_DOMAINS or "youtube.com" in domain:
        return "youtube"
    if domain in TWITTER_DOMAINS:
        return "twitter"
    if domain in REDDIT_DOMAINS:
        return "reddit"
    if domain in ARXIV_DOMAINS:
        return "arxiv"
    return "article"


def is_good_subtitles(text: str | None) -> bool:
    if not text or len(text) < SUBTITLES_MIN_CHARS:
        return False
    lines = [l.strip() for l in text.split("\n") if l.strip()]
    if len(lines) < 5:
        return False
    unique_ratio = len(set(lines)) / len(lines)
    return unique_ratio >= SUBTITLES_MIN_UNIQUE_RATIO


async def triage_youtube(url: str) -> tuple[dict, list[str]]:
    """
    Check YouTube subtitles availability.
    Returns (triage_data, required_capabilities).
    """
    try:
        import yt_dlp
    except ImportError:
        log.warning("yt_dlp not installed, skipping subtitle check")
        return {"subtitles_available": False}, ["heavy_transcription", "summarization"]

    ydl_opts = {
        "quiet": True,
        "no_warnings": True,
        "skip_download": True,
        "writesubtitles": False,
        "writeautomaticsub": False,
        "subtitleslangs": ["ru", "en", "ru-RU"],
        "socket_timeout": 15,
    }

    triage_data: dict = {}

    try:
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=False)

        duration = info.get("duration", 0)
        triage_data["duration_seconds"] = duration
        triage_data["title"] = info.get("title", "")
        triage_data["language"] = info.get("language") or info.get("subtitles", {}) and next(
            iter(info.get("subtitles", {})), None
        )

        # Check manual subtitles first, then auto-generated
        subs = info.get("subtitles", {})
        auto_subs = info.get("automatic_captions", {})

        sub_text = None
        for lang_priority in ["ru", "ru-RU", "en"]:
            if lang_priority in subs or lang_priority in auto_subs:
                sub_text = _extract_subtitle_text(
                    subs.get(lang_priority) or auto_subs.get(lang_priority)
                )
                if sub_text:
                    triage_data["language"] = lang_priority
                    break

        if is_good_subtitles(sub_text):
            triage_data["subtitles_available"] = True
            triage_data["subtitles_text"] = sub_text
            required = ["summarization"]
        else:
            triage_data["subtitles_available"] = False
            # Long videos need heavy transcription
            required = ["heavy_transcription", "summarization"] if duration > 1800 else ["light_transcription", "summarization"]

    except Exception as e:
        log.error("YouTube triage failed for %s: %s", url, e)
        triage_data["subtitles_available"] = False
        triage_data["triage_error"] = str(e)
        required = ["heavy_transcription", "summarization"]

    return triage_data, required


def _extract_subtitle_text(sub_formats: list | None) -> str | None:
    if not sub_formats:
        return None
    # Prefer json3 or vtt format with actual text
    for fmt in sub_formats:
        if fmt.get("ext") in ("json3", "srv3", "vtt", "ttml"):
            # In actual usage, yt_dlp would need to download; here we check availability
            return None  # Text available only after download
    return None


async def triage_article(url: str) -> tuple[dict, list[str]]:
    """Extract article text on VPS."""
    try:
        import trafilatura
        import httpx
    except ImportError:
        return {}, ["summarization"]

    triage_data: dict = {}
    try:
        async with httpx.AsyncClient(timeout=20, follow_redirects=True) as client:
            resp = await client.get(url, headers={"User-Agent": "Mozilla/5.0"})
            resp.raise_for_status()
            html = resp.text

        text = trafilatura.extract(
            html,
            include_comments=False,
            include_tables=True,
            no_fallback=False,
        )
        title = trafilatura.extract_metadata(html)
        if title:
            triage_data["article_title"] = title.title or ""

        if text and len(text) > 200:
            triage_data["article_text"] = text
            triage_data["word_count"] = len(text.split())
        else:
            triage_data["extraction_failed"] = True

    except Exception as e:
        log.error("Article triage failed for %s: %s", url, e)
        triage_data["extraction_failed"] = True
        triage_data["triage_error"] = str(e)

    return triage_data, ["summarization"]


async def triage_arxiv(url: str) -> tuple[dict, list[str]]:
    """Fetch ArXiv abstract via API."""
    triage_data: dict = {}
    try:
        import httpx
        # Extract arxiv ID from URL
        match = re.search(r"arxiv\.org/(abs|pdf)/([0-9]+\.[0-9]+)", url)
        if not match:
            return triage_data, ["summarization"]

        arxiv_id = match.group(2)
        api_url = f"https://export.arxiv.org/api/query?id_list={arxiv_id}"

        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.get(api_url)
            resp.raise_for_status()

        # Simple XML extraction without lxml
        text = resp.text
        title_match = re.search(r"<title>(.*?)</title>", text, re.DOTALL)
        abstract_match = re.search(r"<summary>(.*?)</summary>", text, re.DOTALL)

        if title_match:
            triage_data["article_title"] = title_match.group(1).strip()
        if abstract_match:
            triage_data["article_text"] = abstract_match.group(1).strip()
        triage_data["arxiv_id"] = arxiv_id
        triage_data["pdf_url"] = f"https://arxiv.org/pdf/{arxiv_id}"

    except Exception as e:
        log.error("ArXiv triage failed for %s: %s", url, e)
        triage_data["triage_error"] = str(e)

    return triage_data, ["summarization"]


async def triage_pdf(file_path: str) -> tuple[dict, list[str]]:
    """Try to extract text from PDF; flag for OCR if needed."""
    triage_data: dict = {}
    try:
        import fitz  # PyMuPDF

        full_path = f"media_staging/{file_path}" if not file_path.startswith("media_staging") else file_path
        doc = fitz.open(full_path)
        text_parts = []
        for page in doc:
            text_parts.append(page.get_text())
        text = "\n".join(text_parts).strip()

        if len(text) > 200:
            triage_data["pdf_text"] = text
            triage_data["needs_ocr"] = False
            return triage_data, ["summarization"]
        else:
            triage_data["needs_ocr"] = True
            return triage_data, ["pdf_ocr", "summarization"]

    except ImportError:
        triage_data["needs_ocr"] = True
        return triage_data, ["pdf_ocr", "summarization"]
    except Exception as e:
        log.error("PDF triage failed for %s: %s", file_path, e)
        triage_data["triage_error"] = str(e)
        return triage_data, ["pdf_ocr", "summarization"]


def capabilities_for_media(task_type: str, duration_seconds: float | None = None) -> list[str]:
    """Return required capabilities for audio/video tasks."""
    heavy = duration_seconds and duration_seconds > 1800
    transcription = "heavy_transcription" if heavy else "light_transcription"
    return [transcription, "summarization"]


OFFICE_TYPES = {"word", "powerpoint", "excel", "odf", "rtf", "epub"}


async def triage_office(file_path: str, task_type: str) -> tuple[dict, list[str]]:
    """
    Try to extract text from office documents on VPS.
    Returns (triage_data, required_capabilities).
    """
    triage_data: dict = {"doc_type": task_type}

    try:
        import fitz  # PyMuPDF handles some office formats

        if task_type == "word":
            text = _extract_word(file_path)
        elif task_type == "powerpoint":
            text = _extract_pptx(file_path)
        elif task_type == "excel":
            text = _extract_xlsx(file_path)
        elif task_type in ("odf", "rtf"):
            text = _extract_generic_text(file_path)
        elif task_type == "epub":
            text = _extract_epub(file_path)
        else:
            text = None

        if text and len(text.strip()) > 100:
            triage_data["article_text"] = text
            return triage_data, ["office_document", "summarization"]

    except Exception as e:
        log.error("Office triage failed for %s (%s): %s", file_path, task_type, e)
        triage_data["triage_error"] = str(e)

    # Fallback: worker handles it
    return triage_data, ["office_document", "summarization"]


def _extract_word(file_path: str) -> str | None:
    try:
        from docx import Document
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except ImportError:
        return None


def _extract_pptx(file_path: str) -> str | None:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            if texts:
                slides.append(f"[Слайд {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except ImportError:
        return None


def _extract_xlsx(file_path: str) -> str | None:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
        sheets = []
        for ws in wb.worksheets:
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    rows.append("\t".join(cells))
            if rows:
                sheets.append(f"[Лист: {ws.title}]\n" + "\n".join(rows[:100]))  # cap at 100 rows
        return "\n\n".join(sheets)
    except ImportError:
        return None


def _extract_epub(file_path: str) -> str | None:
    try:
        import ebooklib
        from ebooklib import epub
        from html import unescape
        import re as _re
        book = epub.read_epub(file_path)
        parts = []
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            raw = item.get_content().decode("utf-8", errors="ignore")
            text = unescape(_re.sub(r"<[^>]+>", " ", raw)).strip()
            if text:
                parts.append(text)
        return "\n\n".join(parts)
    except ImportError:
        return None


def _extract_generic_text(file_path: str) -> str | None:
    """Fallback: read as plain text (works for RTF with some noise)."""
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception:
        return None
